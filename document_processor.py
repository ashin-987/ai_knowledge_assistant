"""
Improved Document Processor - Smart Chunking & Multi-Format Support
Supports: PDF, TXT, DOCX, PPTX with semantic chunking
"""

from pathlib import Path
from typing import List, Dict, Optional
import PyPDF2
from datetime import datetime
import hashlib

# Optional imports - graceful degradation if not installed
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

class DocumentProcessor:
    def __init__(self, chunk_size=1000, chunk_overlap=200):
        """
        Initialize the document processor with smart chunking.
        
        Args:
            chunk_size: Target characters per chunk
            chunk_overlap: Characters to overlap between chunks
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        
        # Separators for semantic chunking (in order of preference)
        self.separators = [
            "\n\n\n",  # Multiple line breaks (section boundaries)
            "\n\n",    # Paragraph breaks
            "\n",      # Line breaks
            ". ",      # Sentence boundaries
            "! ",      # Sentence boundaries
            "? ",      # Sentence boundaries
            "; ",      # Clause boundaries
            ", ",      # Phrase boundaries
            " ",       # Word boundaries
            ""         # Character-level fallback
        ]
    
    def chunk_text_semantic(self, text: str, metadata: Optional[Dict] = None) -> List[Dict]:
        """
        Split text into chunks that respect semantic boundaries.
        
        Args:
            text: Text to chunk
            metadata: Optional metadata to attach to chunks
            
        Returns:
            List of chunk dictionaries
        """
        if not text or not text.strip():
            return []
        
        chunks = []
        current_chunk = ""
        
        # Recursive splitting using separators
        splits = self._split_text_recursive(text, self.separators)
        
        for split in splits:
            # If adding this split would exceed chunk_size, save current chunk
            if len(current_chunk) + len(split) > self.chunk_size and current_chunk:
                chunks.append(current_chunk.strip())
                
                # Start new chunk with overlap
                overlap_text = self._get_overlap(current_chunk)
                current_chunk = overlap_text + split
            else:
                current_chunk += split
        
        # Add the last chunk
        if current_chunk.strip():
            chunks.append(current_chunk.strip())
        
        # Create chunk dictionaries with metadata
        chunk_dicts = []
        for i, chunk in enumerate(chunks):
            chunk_dict = {
                'text': chunk,
                'chunk_id': i,
                'metadata': metadata or {}
            }
            chunk_dicts.append(chunk_dict)
        
        return chunk_dicts
    
    def _split_text_recursive(self, text: str, separators: List[str]) -> List[str]:
        """Recursively split text using hierarchical separators."""
        if not separators:
            return [text]
        
        separator = separators[0]
        remaining_separators = separators[1:]
        
        if separator == "":
            # Character-level split as fallback
            return [text[i:i+self.chunk_size] for i in range(0, len(text), self.chunk_size)]
        
        splits = text.split(separator)
        
        # If we got good splits, return them (with separator reattached)
        if len(splits) > 1:
            result = []
            for i, split in enumerate(splits):
                if i < len(splits) - 1:
                    # Reattach separator (except for last split)
                    result.append(split + separator)
                else:
                    result.append(split)
            return [s for s in result if s.strip()]
        
        # If no split happened, try next separator
        return self._split_text_recursive(text, remaining_separators)
    
    def _get_overlap(self, text: str) -> str:
        """Get overlap text from end of chunk."""
        if len(text) <= self.chunk_overlap:
            return text
        return text[-self.chunk_overlap:]
    
    def extract_metadata(self, file_path: str) -> Dict:
        """
        Extract metadata from file.
        
        Args:
            file_path: Path to the file
            
        Returns:
            Dictionary of metadata
        """
        path = Path(file_path)
        metadata = {
            'filename': path.name,
            'file_type': path.suffix.lower(),
            'file_size': path.stat().st_size,
            'processed_date': datetime.now().isoformat()
        }
        
        # PDF-specific metadata
        if path.suffix.lower() == '.pdf':
            try:
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    metadata.update({
                        'page_count': len(pdf_reader.pages),
                        'author': pdf_reader.metadata.get('/Author', '') if pdf_reader.metadata else '',
                        'title': pdf_reader.metadata.get('/Title', '') if pdf_reader.metadata else '',
                        'creation_date': str(pdf_reader.metadata.get('/CreationDate', '')) if pdf_reader.metadata else ''
                    })
            except Exception as e:
                print(f"⚠️ Could not extract PDF metadata: {e}")
        
        return metadata
    
    def load_pdf(self, pdf_path: str) -> tuple[str, Dict]:
        """
        Extract text and metadata from a PDF file.
        
        Returns:
            Tuple of (text, metadata)
        """
        text = ""
        metadata = self.extract_metadata(pdf_path)
        
        try:
            with open(pdf_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                for page_num, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    if page_text:
                        # Add page markers for better context
                        text += f"\n--- Page {page_num + 1} ---\n{page_text}\n"
                
                metadata['pages_processed'] = len(pdf_reader.pages)
                
        except Exception as e:
            print(f"❌ Error reading PDF {pdf_path}: {e}")
            metadata['error'] = str(e)
        
        return text, metadata
    
    def load_txt(self, txt_path: str) -> tuple[str, Dict]:
        """Load a text file."""
        metadata = self.extract_metadata(txt_path)
        
        try:
            with open(txt_path, 'r', encoding='utf-8') as file:
                text = file.read()
            metadata['encoding'] = 'utf-8'
            return text, metadata
        except UnicodeDecodeError:
            # Try different encodings
            for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
                try:
                    with open(txt_path, 'r', encoding=encoding) as file:
                        text = file.read()
                    metadata['encoding'] = encoding
                    return text, metadata
                except:
                    continue
            
            print(f"❌ Error reading text file {txt_path}: encoding issue")
            metadata['error'] = 'encoding_error'
            return "", metadata
        except Exception as e:
            print(f"❌ Error reading text file {txt_path}: {e}")
            metadata['error'] = str(e)
            return "", metadata
    
    def load_docx(self, docx_path: str) -> tuple[str, Dict]:
        """Extract text from Word document."""
        metadata = self.extract_metadata(docx_path)
        
        if not DOCX_AVAILABLE:
            print("⚠️ python-docx not installed. Install with: pip install python-docx")
            metadata['error'] = 'python-docx not installed'
            return "", metadata
        
        try:
            doc = Document(docx_path)
            
            # Extract paragraphs
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
            text = "\n\n".join(paragraphs)
            
            # Extract tables
            tables_text = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(" | ".join(row_data))
                tables_text.append("\n".join(table_data))
            
            if tables_text:
                text += "\n\n=== Tables ===\n\n" + "\n\n".join(tables_text)
            
            metadata['paragraph_count'] = len(paragraphs)
            metadata['table_count'] = len(doc.tables)
            
            return text, metadata
            
        except Exception as e:
            print(f"❌ Error reading DOCX {docx_path}: {e}")
            metadata['error'] = str(e)
            return "", metadata
    
    def load_pptx(self, pptx_path: str) -> tuple[str, Dict]:
        """Extract text from PowerPoint presentation."""
        metadata = self.extract_metadata(pptx_path)
        
        if not PPTX_AVAILABLE:
            print("⚠️ python-pptx not installed. Install with: pip install python-pptx")
            metadata['error'] = 'python-pptx not installed'
            return "", metadata
        
        try:
            prs = Presentation(pptx_path)
            
            slides_text = []
            for i, slide in enumerate(prs.slides):
                slide_text = f"\n--- Slide {i + 1} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                
                slides_text.append(slide_text)
            
            text = "\n".join(slides_text)
            metadata['slide_count'] = len(prs.slides)
            
            return text, metadata
            
        except Exception as e:
            print(f"❌ Error reading PPTX {pptx_path}: {e}")
            metadata['error'] = str(e)
            return "", metadata
    
    def process_file(self, file_path: str) -> List[Dict]:
        """
        Process a single file and return chunks with metadata.
        
        Args:
            file_path: Path to the file
            
        Returns:
            List of chunk dictionaries
        """
        path = Path(file_path)
        suffix = path.suffix.lower()
        
        print(f"📄 Processing: {path.name}")
        
        # Load file based on type
        if suffix == '.pdf':
            text, metadata = self.load_pdf(str(path))
        elif suffix == '.txt':
            text, metadata = self.load_txt(str(path))
        elif suffix == '.docx':
            text, metadata = self.load_docx(str(path))
        elif suffix == '.pptx':
            text, metadata = self.load_pptx(str(path))
        else:
            print(f"⚠️ Unsupported file type: {suffix}")
            return []
        
        if not text.strip():
            print(f"⚠️ No text extracted from {path.name}")
            return []
        
        # Chunk the text semantically
        chunk_dicts = self.chunk_text_semantic(text, metadata)
        
        # Add source to each chunk
        for chunk in chunk_dicts:
            chunk['source'] = path.name
        
        print(f"✅ Created {len(chunk_dicts)} chunks from {path.name}")
        
        return chunk_dicts
    
    def process_directory(self, directory: str) -> List[Dict]:
        """
        Process all supported files in a directory.
        
        Args:
            directory: Path to directory
            
        Returns:
            List of all chunks from all files
        """
        print(f"\n📂 Processing directory: {directory}")
        
        supported_extensions = {'.pdf', '.txt', '.docx', '.pptx'}
        all_documents = []
        
        # Find all supported files
        dir_path = Path(directory)
        files = []
        for ext in supported_extensions:
            files.extend(list(dir_path.rglob(f'*{ext}')))
        
        if not files:
            print("⚠️ No supported files found!")
            print(f"Supported formats: {', '.join(supported_extensions)}")
            return []
        
        print(f"📚 Found {len(files)} files to process")
        
        # Process each file
        for file_path in files:
            chunks = self.process_file(str(file_path))
            all_documents.extend(chunks)
        
        print(f"\n✅ Total chunks created: {len(all_documents)}")
        print(f"📊 Files processed: {len(files)}")
        
        return all_documents
    
    def generate_summary(self, text: str, max_length: int = 500) -> str:
        """
        Generate a simple extractive summary.
        
        Args:
            text: Text to summarize
            max_length: Maximum summary length
            
        Returns:
            Summary text
        """
        # Simple extractive summary: take first few sentences
        sentences = text.split('. ')
        summary = ""
        
        for sentence in sentences:
            if len(summary) + len(sentence) < max_length:
                summary += sentence + '. '
            else:
                break
        
        return summary.strip()
